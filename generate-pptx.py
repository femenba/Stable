"""
Generate stable. investor pitch deck as a PPTX file.
Run: python3 generate-pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree
import copy

# ── Widescreen 16:9 ──────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

# ── Colour palette ───────────────────────────────────────────
BG_DARK     = RGBColor(0x0a, 0x08, 0x14)
BG_CARD     = RGBColor(0x14, 0x10, 0x2a)
PURPLE_L    = RGBColor(0xa7, 0x8b, 0xfa)   # #a78bfa
PURPLE_D    = RGBColor(0x63, 0x66, 0xf1)   # #6366f1
VIOLET      = RGBColor(0xa8, 0x55, 0xf7)   # #a855f7
GOLD        = RGBColor(0xfb, 0xbf, 0x24)   # #fbbf24
GREEN       = RGBColor(0x4a, 0xde, 0x80)   # #4ade80
RED         = RGBColor(0xf8, 0x71, 0x71)   # #f87171
WHITE       = RGBColor(0xff, 0xff, 0xff)
WHITE_DIM   = RGBColor(0xbb, 0xbb, 0xcc)
WHITE_MUTED = RGBColor(0x66, 0x60, 0x80)
CARD_BORDER = RGBColor(0x2a, 0x20, 0x50)

def rgb_hex(r):
    return "{:02X}{:02X}{:02X}".format(r[0], r[1], r[2])

def add_slide(layout_index=6):
    layout = prs.slide_layouts[layout_index]  # blank
    return prs.slides.add_slide(layout)

def fill_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(0)):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, left, top, width, height,
             font_size=Pt(14), bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_label(slide, text, left, top, width=Inches(10)):
    """Small uppercase label in purple"""
    add_text(slide, text.upper(), left, top, width, Inches(0.3),
             font_size=Pt(9), bold=True, color=PURPLE_L)

def add_title(slide, text, left, top, width=Inches(11), size=Pt(36)):
    add_text(slide, text, left, top, width, Inches(1.5),
             font_size=size, bold=True, color=WHITE)

def add_divider(slide, left, top, width=Inches(0.6)):
    add_rect(slide, left, top, width, Pt(4), fill_color=PURPLE_D)

def add_card(slide, left, top, width, height, color=CARD_BORDER, fill=None):
    fc = fill or RGBColor(0x16, 0x12, 0x2e)
    r = add_rect(slide, left, top, width, height,
                 fill_color=fc, line_color=color, line_width=Pt(1))
    return r

def slide_num_label(slide, n, total=11):
    add_text(slide, f"{n:02d} / {total:02d}",
             W - Inches(1.4), Pt(16), Inches(1.2), Inches(0.3),
             font_size=Pt(9), color=WHITE_MUTED, align=PP_ALIGN.RIGHT)

# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ═══════════════════════════════════════════════════════════════
s = add_slide()
fill_bg(s, BG_DARK)
slide_num_label(s, 1)

# Gradient-ish background rectangle
add_rect(s, 0, 0, W, H*0.6, fill_color=RGBColor(0x1e, 0x12, 0x60))
add_rect(s, 0, H*0.4, W, H*0.6, fill_color=BG_DARK)

# Logo wordmark
add_text(s, "stable.", Inches(0), Inches(1.8), W, Inches(1.6),
         font_size=Pt(80), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Tagline
add_text(s, "ADHD & FOCUS PRODUCTIVITY PLATFORM",
         Inches(0), Inches(3.3), W, Inches(0.5),
         font_size=Pt(11), bold=True, color=PURPLE_L, align=PP_ALIGN.CENTER)

# CTA pill (rectangle + text)
pill = add_rect(s, Inches(4.5), Inches(4.2), Inches(4.3), Inches(0.6),
                fill_color=PURPLE_D, line_color=VIOLET)
add_text(s, "INVESTMENT OPPORTUNITY · 2026",
         Inches(4.5), Inches(4.2), Inches(4.3), Inches(0.6),
         font_size=Pt(11), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Confidential
add_text(s, "CONFIDENTIAL · FOR QUALIFIED BUYERS ONLY",
         Inches(0), H - Inches(0.5), W, Inches(0.35),
         font_size=Pt(8), color=WHITE_MUTED, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ═══════════════════════════════════════════════════════════════
s = add_slide()
fill_bg(s, BG_DARK)
slide_num_label(s, 2)

add_label(s, "The Problem", Inches(0.7), Inches(0.5))
add_title(s, "Most productivity apps fail ADHD brains.", Inches(0.7), Inches(0.85), size=Pt(34))
add_divider(s, Inches(0.7), Inches(1.85))

stats = [
    ("366M",  "Adults worldwide with ADHD",
     "The most underserved segment in the productivity market."),
    ("~$0",   "Apps built specifically for them",
     "Existing tools (Todoist, Notion, Asana) overwhelm ADHD users with complexity."),
    ("71%",   "Abandon general productivity apps within 30 days",
     "Too many features. Too much friction. Not built for them."),
]

for i, (num, lbl, body) in enumerate(stats):
    cx = Inches(0.7) + i * Inches(4.1)
    add_card(s, cx, Inches(2.1), Inches(3.9), Inches(4.0))
    add_text(s, num, cx + Inches(0.2), Inches(2.3), Inches(3.5), Inches(1.2),
             font_size=Pt(48), bold=True, color=RGBColor(0xf8,0x71,0x71))
    add_text(s, lbl, cx + Inches(0.2), Inches(3.4), Inches(3.5), Inches(0.5),
             font_size=Pt(10), bold=True, color=PURPLE_L)
    add_text(s, body, cx + Inches(0.2), Inches(3.9), Inches(3.5), Inches(1.2),
             font_size=Pt(11), color=WHITE_MUTED)

# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — THE SOLUTION
# ═══════════════════════════════════════════════════════════════
s = add_slide()
fill_bg(s, BG_DARK)
slide_num_label(s, 3)

add_label(s, "The Solution", Inches(0.7), Inches(0.5))
add_title(s, "stable. — the focus app\nbuilt for ADHD adults.", Inches(0.7), Inches(0.85), size=Pt(30))
add_divider(s, Inches(0.7), Inches(1.95))

features = [
    "Three tasks. That's it. — radical simplicity ADHD brains love",
    "AI insight card — personalised daily coaching, no setup needed",
    "Task categories — Work, Personal, Family, Health, Other",
    "Focus sessions — built-in Pomodoro-style timer + history",
    "Smart reminders — snooze, dismiss, recurring options",
    "Light + dark mode — polished, professional UI",
    "Web + Mobile (Expo) — works on every device",
]
for i, f in enumerate(features):
    add_rect(s, Inches(0.7), Inches(2.2) + i*Inches(0.58),
             Inches(0.12), Inches(0.12),
             fill_color=PURPLE_D)
    add_text(s, f, Inches(1.0), Inches(2.15) + i*Inches(0.58),
             Inches(7.0), Inches(0.5), font_size=Pt(13), color=WHITE_DIM)

# Phone mockup (simplified)
phone_l = Inches(9.2)
phone_t = Inches(0.7)
phone_w = Inches(3.4)
phone_h = Inches(6.3)
add_rect(s, phone_l, phone_t, phone_w, phone_h,
         fill_color=RGBColor(0x12,0x0e,0x24),
         line_color=RGBColor(0x2a,0x1f,0x60), line_width=Pt(1.5))
# Header gradient block
add_rect(s, phone_l, phone_t, phone_w, Inches(1.6),
         fill_color=RGBColor(0x4f,0x3a,0xff))
add_text(s, "TODAY'S FOCUS", phone_l+Inches(0.2), phone_t+Inches(0.15),
         Inches(3.0), Inches(0.3), font_size=Pt(7), color=WHITE_MUTED, bold=True)
add_text(s, "Three things.\nThat's it.", phone_l+Inches(0.2), phone_t+Inches(0.45),
         Inches(3.0), Inches(0.9), font_size=Pt(16), bold=True, color=WHITE)
# Task cards
task_colors = [PURPLE_D, VIOLET, RGBColor(0xec,0x48,0x99)]
task_names  = ["Review Q2 budget", "Book check-up", "30 min walk"]
task_tags   = ["Work", "Family", "Health"]
for i in range(3):
    cy = phone_t + Inches(1.75) + i * Inches(1.1)
    add_rect(s, phone_l + Inches(0.15), cy, phone_w - Inches(0.3), Inches(0.95),
             fill_color=RGBColor(0x1a,0x15,0x35), line_color=CARD_BORDER, line_width=Pt(0.5))
    add_rect(s, phone_l + Inches(0.15), cy, Inches(0.06), Inches(0.95),
             fill_color=task_colors[i])
    add_text(s, task_names[i], phone_l+Inches(0.3), cy+Inches(0.1),
             Inches(2.8), Inches(0.4), font_size=Pt(11), bold=True, color=WHITE)
    add_text(s, task_tags[i], phone_l+Inches(0.3), cy+Inches(0.5),
             Inches(1.0), Inches(0.3), font_size=Pt(8), color=PURPLE_L)

# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — MARKET OPPORTUNITY
# ═══════════════════════════════════════════════════════════════
s = add_slide()
fill_bg(s, BG_DARK)
slide_num_label(s, 4)

add_label(s, "Market Opportunity", Inches(0.7), Inches(0.5))
add_title(s, "A $25 billion market with no clear winner.", Inches(0.7), Inches(0.85), size=Pt(30))
add_divider(s, Inches(0.7), Inches(1.75))

mkt = [
    ("366M",  "Adults\nwith ADHD globally",   PURPLE_D),
    ("$25B",  "ADHD management\nmarket by 2030", VIOLET),
    ("$85B",  "Productivity app\nmarket by 2027", PURPLE_L),
    ("~10%",  "Of adults have\nADHD-adjacent traits", GREEN),
]
for i, (num, lbl, col) in enumerate(mkt):
    cx = Inches(0.7) + i * Inches(3.15)
    add_card(s, cx, Inches(2.0), Inches(3.0), Inches(2.6),
             color=col, fill=RGBColor(0x16,0x12,0x2e))
    add_text(s, num, cx+Inches(0.2), Inches(2.1), Inches(2.6), Inches(0.9),
             font_size=Pt(36), bold=True, color=col)
    add_text(s, lbl, cx+Inches(0.2), Inches(2.95), Inches(2.6), Inches(0.7),
             font_size=Pt(11), color=WHITE_MUTED)

# Bottom insight box
add_rect(s, Inches(0.7), Inches(4.8), Inches(12.0), Inches(1.6),
         fill_color=RGBColor(0x16,0x12,0x35),
         line_color=RGBColor(0x36,0x30,0x6a), line_width=Pt(1))
add_text(s, "No dominant player in ADHD-specific productivity.\n"
            "Closest competitors (Focusmate, Inflow) are single-feature or coaching-only.\n"
            "stable. is the only AI-powered, full-stack focus OS for ADHD adults.",
         Inches(1.0), Inches(4.9), Inches(11.5), Inches(1.4),
         font_size=Pt(12), color=WHITE_DIM)

# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — TECH & STATUS
# ═══════════════════════════════════════════════════════════════
s = add_slide()
fill_bg(s, BG_DARK)
slide_num_label(s, 5)

add_label(s, "Product Status", Inches(0.7), Inches(0.5))
add_title(s, "Built, deployed, ready to scale.", Inches(0.7), Inches(0.85), size=Pt(34))
add_divider(s, Inches(0.7), Inches(1.75))

# Web app card
add_card(s, Inches(0.7), Inches(2.0), Inches(5.8), Inches(1.5),
         color=GREEN, fill=RGBColor(0x0c,0x1f,0x15))
add_text(s, "✅  Web App — LIVE", Inches(0.9), Inches(2.1), Inches(5.4), Inches(0.45),
         font_size=Pt(14), bold=True, color=GREEN)
add_text(s, "Next.js 15 deployed on Vercel. Full auth, database, real-time API. All 4 core features complete.",
         Inches(0.9), Inches(2.5), Inches(5.4), Inches(0.7),
         font_size=Pt(11), color=WHITE_MUTED)

# Mobile card
add_card(s, Inches(0.7), Inches(3.65), Inches(5.8), Inches(1.5),
         color=GOLD, fill=RGBColor(0x1c,0x16,0x07))
add_text(s, "🚧  Mobile App — In Development", Inches(0.9), Inches(3.75), Inches(5.4), Inches(0.45),
         font_size=Pt(14), bold=True, color=GOLD)
add_text(s, "React Native (Expo). Shares same design system and API. iOS + Android support.",
         Inches(0.9), Inches(4.15), Inches(5.4), Inches(0.7),
         font_size=Pt(11), color=WHITE_MUTED)

# Tech stack list
add_text(s, "TECH STACK", Inches(7.2), Inches(2.0), Inches(5.5), Inches(0.3),
         font_size=Pt(9), bold=True, color=PURPLE_L)
tech = [
    "Next.js 15 + TypeScript — modern scalable frontend",
    "tRPC + React Query — fully type-safe API",
    "Supabase (PostgreSQL) — managed, instant scaling",
    "Clerk — enterprise-grade authentication",
    "Vercel — zero-config global deployment",
    "Expo (React Native) — iOS + Android from one codebase",
    "pnpm monorepo — shared code across web and mobile",
]
for i, t in enumerate(tech):
    add_rect(s, Inches(7.2), Inches(2.45)+i*Inches(0.55), Inches(0.1), Inches(0.1),
             fill_color=PURPLE_D)
    add_text(s, t, Inches(7.45), Inches(2.42)+i*Inches(0.55),
             Inches(5.5), Inches(0.45), font_size=Pt(12), color=WHITE_DIM)

# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — PRICING
# ═══════════════════════════════════════════════════════════════
s = add_slide()
fill_bg(s, BG_DARK)
slide_num_label(s, 6)

add_label(s, "Monetisation Model", Inches(0.7), Inches(0.5))
add_title(s, "Four tiers. Freemium to enterprise.", Inches(0.7), Inches(0.85), size=Pt(32))
add_divider(s, Inches(0.7), Inches(1.75))

tiers = [
    ("FREE",    "£0",    "/forever",     WHITE_MUTED, RGBColor(0x2a,0x20,0x50), [
        "3 active tasks", "1 focus session/day", "Basic reminders", "Light mode"]),
    ("PRO",     "£7.99", "/month",       PURPLE_L,    RGBColor(0x1e,0x1a,0x55), [
        "Unlimited tasks", "All categories", "Full AI insights", "Focus analytics",
        "Dark + light mode", "Priority support"]),
    ("FAMILY",  "£13.99","/month",       VIOLET,      RGBColor(0x1a,0x12,0x40), [
        "Everything in Pro", "4 user accounts", "Family reminders", "Shared tasks"]),
    ("TEAMS",   "£6.99", "/user/month",  GOLD,        RGBColor(0x1c,0x16,0x07), [
        "Everything in Pro", "Team dashboard", "Admin controls",
        "Invoiced billing", "ADHD onboarding"]),
]

for i, (name, price, period, col, fill_c, feats) in enumerate(tiers):
    cx = Inches(0.55) + i * Inches(3.2)
    add_card(s, cx, Inches(2.05), Inches(3.05), Inches(4.9),
             color=col, fill=fill_c)
    add_text(s, name, cx+Inches(0.2), Inches(2.15), Inches(2.65), Inches(0.35),
             font_size=Pt(9), bold=True, color=col)
    add_text(s, price, cx+Inches(0.2), Inches(2.5), Inches(2.65), Inches(0.7),
             font_size=Pt(32), bold=True, color=WHITE)
    add_text(s, period, cx+Inches(0.2), Inches(3.15), Inches(2.65), Inches(0.3),
             font_size=Pt(10), color=WHITE_MUTED)
    for j, f in enumerate(feats):
        add_text(s, "✓  " + f, cx+Inches(0.2), Inches(3.55)+j*Inches(0.46),
                 Inches(2.65), Inches(0.4), font_size=Pt(10), color=WHITE_DIM)

# Most popular badge
add_rect(s, Inches(3.75), Inches(1.9), Inches(1.6), Inches(0.28),
         fill_color=PURPLE_D, line_color=VIOLET)
add_text(s, "MOST POPULAR", Inches(3.75), Inches(1.88), Inches(1.6), Inches(0.28),
         font_size=Pt(7), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — REVENUE PROJECTIONS
# ═══════════════════════════════════════════════════════════════
s = add_slide()
fill_bg(s, BG_DARK)
slide_num_label(s, 7)

add_label(s, "Revenue Projections", Inches(0.7), Inches(0.5))
add_title(s, "£1M+ ARR reachable within 36 months.", Inches(0.7), Inches(0.85), size=Pt(32))
add_divider(s, Inches(0.7), Inches(1.75))

# Table header
headers = ["Milestone", "Paid Users", "Avg Revenue", "Annual Revenue"]
col_x   = [Inches(0.7), Inches(3.2), Inches(5.8), Inches(8.6)]
col_w   = [Inches(2.4), Inches(2.4), Inches(2.6), Inches(2.8)]
for i, h in enumerate(headers):
    add_text(s, h.upper(), col_x[i], Inches(2.15), col_w[i], Inches(0.35),
             font_size=Pt(9), bold=True, color=WHITE_MUTED)

add_rect(s, Inches(0.7), Inches(2.5), Inches(11.9), Pt(1.5), fill_color=RGBColor(0x2a,0x20,0x50))

rows = [
    ("Month 6",  "200",    "£7.99 / mo",  "£19,200"),
    ("Year 1",   "600",    "£7.99 / mo",  "£57,500"),
    ("Year 2",   "3,000",  "£8.50 / mo",  "£306,000"),
    ("Year 3",   "10,000", "£8.99 / mo",  "£1,080,000"),
]
for ri, (m, u, a, r) in enumerate(rows):
    ry = Inches(2.65) + ri * Inches(0.8)
    bg = RGBColor(0x12,0x0e,0x22) if ri % 2 == 0 else RGBColor(0x16,0x12,0x2c)
    add_rect(s, Inches(0.7), ry, Inches(11.9), Inches(0.75), fill_color=bg)
    vals = [m, u, a, r]
    for ci, v in enumerate(vals):
        is_rev = ci == 3
        add_text(s, v, col_x[ci]+Inches(0.1), ry+Inches(0.15), col_w[ci], Inches(0.5),
                 font_size=Pt(14) if ri == 3 else Pt(13),
                 bold=(ri == 3 or is_rev),
                 color=PURPLE_L if is_rev else (WHITE if ri == 3 else WHITE_DIM))

# Side stats
stats_r = [
    ("📈", "Growth via organic channels\n(App Store SEO, ADHD communities)"),
    ("💰", "~95% gross margin at scale\n(SaaS infra ~£400/mo at 10k users)"),
    ("♻️", "Target <4% monthly churn\n(users who find a system, stick with it)"),
]
for i, (ic, txt) in enumerate(stats_r):
    add_card(s, Inches(11.85) - Inches(0.4), Inches(2.0)+i*Inches(1.65), Inches(1.6), Inches(1.5))
    # These won't render as emoji in PPTX reliably, use text only
    add_text(s, txt, Inches(11.5), Inches(2.1)+i*Inches(1.65), Inches(1.6), Inches(1.3),
             font_size=Pt(9), color=WHITE_MUTED)

# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — COMPETITIVE LANDSCAPE
# ═══════════════════════════════════════════════════════════════
s = add_slide()
fill_bg(s, BG_DARK)
slide_num_label(s, 8)

add_label(s, "Competitive Landscape", Inches(0.7), Inches(0.5))
add_title(s, "No one owns this space yet.", Inches(0.7), Inches(0.85), size=Pt(34))
add_divider(s, Inches(0.7), Inches(1.75))

comps = [
    ("Todoist / Any.do",
     "General productivity. Complex, feature-heavy. No ADHD focus. No AI coaching. Not mobile-first emotional design.",
     "Not built for ADHD"),
    ("Focusmate",
     "Body-doubling sessions. Single feature. No task management, no AI, no reminders. Different use case entirely.",
     "Single-feature only"),
    ("Inflow / Numo",
     "ADHD coaching apps. Heavy on content, light on execution tools. Expensive (£40–60/mo). No task + focus integration.",
     "Coaching only, too costly"),
]
for i, (name, body, verdict) in enumerate(comps):
    cx = Inches(0.7) + i * Inches(4.1)
    add_card(s, cx, Inches(2.05), Inches(3.9), Inches(3.0))
    add_text(s, name, cx+Inches(0.2), Inches(2.2), Inches(3.5), Inches(0.4),
             font_size=Pt(14), bold=True, color=WHITE)
    add_text(s, body, cx+Inches(0.2), Inches(2.65), Inches(3.5), Inches(1.4),
             font_size=Pt(11), color=WHITE_MUTED)
    add_text(s, "✕  " + verdict, cx+Inches(0.2), Inches(4.5), Inches(3.5), Inches(0.35),
             font_size=Pt(10), bold=True, color=RED)

# Bottom winner box
add_rect(s, Inches(0.7), Inches(5.35), Inches(12.0), Inches(1.5),
         fill_color=RGBColor(0x1a,0x16,0x42),
         line_color=PURPLE_D, line_width=Pt(1.5))
add_text(s, "stable. — the only AI-powered ADHD focus OS",
         Inches(1.0), Inches(5.45), Inches(11.5), Inches(0.4),
         font_size=Pt(15), bold=True, color=WHITE)
add_text(s, "Full task management + AI daily coaching + focus sessions + reminders. One app. One price. Built for ADHD brains.",
         Inches(1.0), Inches(5.85), Inches(11.5), Inches(0.75),
         font_size=Pt(12), color=WHITE_MUTED)

# ═══════════════════════════════════════════════════════════════
# SLIDE 9 — WHAT'S INCLUDED
# ═══════════════════════════════════════════════════════════════
s = add_slide()
fill_bg(s, BG_DARK)
slide_num_label(s, 9)

add_label(s, "What You're Buying", Inches(0.7), Inches(0.5))
add_title(s, "Everything. Turnkey. Ready to launch.", Inches(0.7), Inches(0.85), size=Pt(32))
add_divider(s, Inches(0.7), Inches(1.75))

items_l = [
    "Full source code (pnpm monorepo — web + mobile)",
    "Web app deployed on Vercel (production-ready)",
    "Supabase database with full schema + migrations",
    "Clerk authentication (user management, billing-ready)",
    "stable. brand assets — logo, wordmark, colour system",
    "Design system + all approved mockups (web + mobile)",
]
items_r = [
    "Expo mobile app codebase (in development, handed over)",
    "All 4 features live: Tasks, Focus, Reminders, Dashboard",
    "tRPC API — fully typed, easy to extend",
    "Stripe-ready billing integration points",
    "30-day handover support from original developer",
    "Full documentation and implementation plans",
]

for i, item in enumerate(items_l):
    ry = Inches(2.15) + i * Inches(0.72)
    add_rect(s, Inches(0.7), ry+Inches(0.05), Inches(0.24), Inches(0.24),
             fill_color=RGBColor(0x22,0x1a,0x50), line_color=PURPLE_D, line_width=Pt(0.5))
    add_text(s, "✓", Inches(0.7), ry+Inches(0.03), Inches(0.24), Inches(0.25),
             font_size=Pt(9), bold=True, color=PURPLE_L, align=PP_ALIGN.CENTER)
    add_text(s, item, Inches(1.05), ry, Inches(5.5), Inches(0.6),
             font_size=Pt(13), color=WHITE_DIM)

for i, item in enumerate(items_r):
    ry = Inches(2.15) + i * Inches(0.72)
    add_rect(s, Inches(7.0), ry+Inches(0.05), Inches(0.24), Inches(0.24),
             fill_color=RGBColor(0x22,0x1a,0x50), line_color=PURPLE_D, line_width=Pt(0.5))
    add_text(s, "✓", Inches(7.0), ry+Inches(0.03), Inches(0.24), Inches(0.25),
             font_size=Pt(9), bold=True, color=PURPLE_L, align=PP_ALIGN.CENTER)
    add_text(s, item, Inches(7.35), ry, Inches(5.5), Inches(0.6),
             font_size=Pt(13), color=WHITE_DIM)

# ═══════════════════════════════════════════════════════════════
# SLIDE 10 — ASKING PRICE
# ═══════════════════════════════════════════════════════════════
s = add_slide()
fill_bg(s, BG_DARK)
slide_num_label(s, 10)

add_label(s, "Investment Ask", Inches(0.7), Inches(0.5))
add_title(s, "Own a platform with £1M+ ARR potential.", Inches(0.7), Inches(0.85), size=Pt(30))
add_divider(s, Inches(0.7), Inches(1.75))

# Big price
add_text(s, "£65,000", Inches(0.7), Inches(1.95), Inches(6.5), Inches(1.5),
         font_size=Pt(72), bold=True, color=GOLD)
add_text(s, "Asking price · one-time payment",
         Inches(0.7), Inches(3.4), Inches(6.0), Inches(0.4),
         font_size=Pt(13), color=WHITE_MUTED)

# Why it's a bargain
add_card(s, Inches(0.7), Inches(3.95), Inches(5.8), Inches(1.3),
         color=GOLD, fill=RGBColor(0x1c,0x16,0x07))
add_text(s, "Why £65k is a bargain",
         Inches(0.9), Inches(4.05), Inches(5.4), Inches(0.4),
         font_size=Pt(13), bold=True, color=GOLD)
add_text(s, "Year 3 revenue projection of £1.08M ARR = 16× return on investment.\n"
            "Equivalent to hiring one developer for 8 months — except you get a finished product.",
         Inches(0.9), Inches(4.45), Inches(5.4), Inches(0.7),
         font_size=Pt(11), color=WHITE_MUTED)

# Payment options
add_card(s, Inches(0.7), Inches(5.4), Inches(5.8), Inches(1.0),
         color=CARD_BORDER)
add_text(s, "Payment options",
         Inches(0.9), Inches(5.5), Inches(5.4), Inches(0.35),
         font_size=Pt(12), bold=True, color=WHITE)
add_text(s, "Full payment  ·  or 50% upfront + 50% at 6-month milestone (£5k MRR)  ·  Open to discuss",
         Inches(0.9), Inches(5.85), Inches(5.4), Inches(0.45),
         font_size=Pt(11), color=WHITE_MUTED)

# Stats column
stat_items = [
    ("£1.08M",  "Projected Year 3 ARR",     PURPLE_L),
    ("~95%",    "Gross margin at scale",     GREEN),
    ("16×",     "ROI on £65k by Year 3",     GOLD),
    ("0",       "Technical debt",            WHITE_MUTED),
]
for i, (num, lbl, col) in enumerate(stat_items):
    cx = Inches(7.4) + (i % 2) * Inches(2.95)
    ry = Inches(2.05) + (i // 2) * Inches(2.1)
    add_card(s, cx, ry, Inches(2.7), Inches(1.8), color=col,
             fill=RGBColor(0x16,0x12,0x2e))
    add_text(s, num, cx+Inches(0.2), ry+Inches(0.15), Inches(2.3), Inches(0.8),
             font_size=Pt(32), bold=True, color=col)
    add_text(s, lbl, cx+Inches(0.2), ry+Inches(0.95), Inches(2.3), Inches(0.6),
             font_size=Pt(10), color=WHITE_MUTED)

# ═══════════════════════════════════════════════════════════════
# SLIDE 11 — CTA
# ═══════════════════════════════════════════════════════════════
s = add_slide()
fill_bg(s, BG_DARK)
slide_num_label(s, 11)

# Background gradient feel
add_rect(s, 0, 0, W, H*0.55, fill_color=RGBColor(0x1e,0x12,0x60))

# Logo
add_text(s, "stable.", Inches(0), Inches(0.6), W, Inches(0.8),
         font_size=Pt(28), bold=True, color=PURPLE_L, align=PP_ALIGN.CENTER)

# Headline
add_text(s, "Ready to own this market?",
         Inches(1.0), Inches(1.5), Inches(11.3), Inches(1.2),
         font_size=Pt(46), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Subtext
add_text(s, "stable. is a finished product in a massive underserved market.\n"
            "The team is here for a smooth handover. First serious buyer wins.",
         Inches(2.0), Inches(2.9), Inches(9.3), Inches(0.85),
         font_size=Pt(15), color=WHITE_MUTED, align=PP_ALIGN.CENTER)

# CTA pill
add_rect(s, Inches(3.5), Inches(4.0), Inches(6.3), Inches(0.75),
         fill_color=PURPLE_D, line_color=VIOLET)
add_text(s, "Get in touch to start due diligence  →",
         Inches(3.5), Inches(4.0), Inches(6.3), Inches(0.75),
         font_size=Pt(14), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Confidential
add_text(s, "CONFIDENTIAL · NOT FOR DISTRIBUTION · stable. · 2026",
         Inches(0), H - Inches(0.45), W, Inches(0.35),
         font_size=Pt(8), color=WHITE_MUTED, align=PP_ALIGN.CENTER)

# ── Save ─────────────────────────────────────────────────────
output = "/Users/femenba/stable/stable-pitch-deck.pptx"
prs.save(output)
print(f"Saved: {output}")
